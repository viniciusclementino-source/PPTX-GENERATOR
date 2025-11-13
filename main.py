# -*- coding: utf-8 -*-
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
import os

# --- Funções Core (Portadas do Colab) ---

def converter_coordenadas_16_9(coord_texto):
    largura_total = 13.33  # polegadas - Ajustado para 16:9
    altura_total = 7.5     # polegadas - Ajustado para 16:9

    try:
        x1y1, x2y2 = coord_texto.split(";")
        x1, y1 = map(float, x1y1.strip().split("-"))
        x2, y2 = map(float, x2y2.strip().split("-"))

        x1, x2 = min(x1, x2), max(x1, x2)
        y1, y2 = min(y1, y2), max(y1, y2)

        # Verifica se as coordenadas estão dentro do limite 0-16 para x e 0-9 para y
        if not (0 <= x1 <= 16 and 0 <= x2 <= 16 and 0 <= y1 <= 9 and 0 <= y2 <= 9):
            raise ValueError("Coordenadas fora do intervalo 16x9.")

        left = Inches(largura_total * (x1 / 16))
        top = Inches(altura_total * (y1 / 9))
        width = Inches(largura_total * ((x2 - x1) / 16))
        height = Inches(altura_total * ((y2 - y1) / 9))

        return left, top, width, height
    except Exception as e:
        messagebox.showerror("Erro de Coordenada", f"Erro ao converter coordenadas '{coord_texto}': {e}")
        return None, None, None, None
def adicionar_caixa_texto(slide, texto, left, top, width, height, font_size=24):
    if None in [left, top, width, height]: return # Aborta se coordenadas falharam
    try:
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        # Limpar parágrafos padrão que podem vir com o textbox
        while len(text_frame.paragraphs) > 1:
             p = text_frame.paragraphs[-1]
             sp = p._p
             sp.getparent().remove(sp)
        p = text_frame.paragraphs[0]
        p.clear() # Limpa qualquer formatação/texto existente no primeiro parágrafo
        run = p.add_run()
        run.text = texto
        run.font.size = Pt(font_size)
    except Exception as e:
        messagebox.showerror("Erro ao Adicionar Texto", f"Não foi possível adicionar a caixa de texto para o texto '{texto[:30]}...': {e}")

def adicionar_imagem(slide, imagem_bytes, left, top, width, height):
    if None in [left, top, width, height]: return # Aborta se coordenadas falharam
    try:
        image_stream = BytesIO(imagem_bytes)
        # Adiciona imagem, tenta garantir que fique por cima de outros elementos se necessário
        picture = slide.shapes.add_picture(image_stream, left, top, width=width, height=height)
        # Mover para o topo da ordem Z (experimental, pode não ser necessário/funcionar como esperado)
        # slide.shapes._spTree.insert(2, picture._element) # O índice 2 é geralmente o início dos shapes
    except Exception as e:
        messagebox.showerror("Erro ao Adicionar Imagem", f"Não foi possível adicionar a imagem: {e}")

# --- Funções de Layout (Portadas e Adaptadas) ---

def layout_img_left_custom(slide, titulo, texto, imagem_bytes, slide_width, slide_height):
    left_txt, top_txt, width_txt, height_txt = converter_coordenadas_16_9("6-4; 15-8")
    left_tit, top_tit, width_tit, height_tit = converter_coordenadas_16_9("6-4; 15-3")
    left_img, top_img, width_img, height_img = converter_coordenadas_16_9("0-0; 6-9")
    adicionar_imagem(slide, imagem_bytes, left_img, top_img, width_img, height_img)
    adicionar_caixa_texto(slide, titulo, left_tit, top_tit, width_tit, height_tit, font_size=36)
    adicionar_caixa_texto(slide, texto, left_txt, top_txt, width_txt, height_txt)

def layout_img_top_custom(slide, titulo, texto, imagem_bytes, slide_width, slide_height):
    left_img, top_img, width_img, height_img = converter_coordenadas_16_9("0-0; 16-4")
    left_tit, top_tit, width_tit, height_tit = converter_coordenadas_16_9("1.5-4.2; 14.5-5.6")
    left_txt, top_txt, width_txt, height_txt = converter_coordenadas_16_9("1-6.6; 15-9")
    adicionar_imagem(slide, imagem_bytes, left_img, top_img, width_img, height_img)
    adicionar_caixa_texto(slide, titulo, left_tit, top_tit, width_tit, height_tit, font_size=36)
    adicionar_caixa_texto(slide, texto, left_txt, top_txt, width_txt, height_txt)

def layout_img_right_custom(slide, titulo, texto, imagem_bytes, slide_width, slide_height):
    left_img, top_img, width_img, height_img = converter_coordenadas_16_9("10-0; 16-9")
    left_tit, top_tit, width_tit, height_tit = converter_coordenadas_16_9("1-1; 10-3")
    left_txt, top_txt, width_txt, height_txt = converter_coordenadas_16_9("1-3; 10-8")
    adicionar_imagem(slide, imagem_bytes, left_img, top_img, width_img, height_img)
    adicionar_caixa_texto(slide, titulo, left_tit, top_tit, width_tit, height_tit, font_size=36)
    adicionar_caixa_texto(slide, texto, left_txt, top_txt, width_txt, height_txt)

def layout_img2(slide, titulo, texto, imagem_bytes, slide_width, slide_height):
    left_tit, top_tit, width_tit, height_tit = converter_coordenadas_16_9("1-1; 15-2.5")
    left_txt, top_txt, width_txt, height_txt = converter_coordenadas_16_9("1-2.6; 15-4.5")
    left_img1, top_img1, width_img1, height_img1 = converter_coordenadas_16_9("1-4.6; 7.5-8.5")
    left_img2, top_img2, width_img2, height_img2 = converter_coordenadas_16_9("8.5-4.6; 15-8.5")
    adicionar_caixa_texto(slide, titulo, left_tit, top_tit, width_tit, height_tit, font_size=36)
    adicionar_caixa_texto(slide, texto, left_txt, top_txt, width_txt, height_txt)
    # Usa a mesma imagem para ambas as posições, como no original
    adicionar_imagem(slide, imagem_bytes, left_img1, top_img1, width_img1, height_img1)
    adicionar_imagem(slide, imagem_bytes, left_img2, top_img2, width_img2, height_img2)

def layout_img6(slide, titulo, *legendas, imagens_bytes, slide_width, slide_height):
    # Coordenadas para as 6 imagens e 6 legendas correspondentes
    coords = [
        # Imagem 1, Legenda 1
        ("0.303-1.303; 5.303-4.303", "0.303-4.403; 5.303-5.103"),
        # Imagem 2, Legenda 2
        ("5.503-1.303; 10.503-4.303", "5.503-4.403; 10.503-5.103"),
        # Imagem 3, Legenda 3
        ("10.703-1.303; 15.703-4.303", "10.703-4.403; 15.703-5.103"),
        # Imagem 4, Legenda 4
        ("0.303-5.206; 5.303-8.206", "0.303-8.306; 5.303-8.906"),
        # Imagem 5, Legenda 5
        ("5.503-5.206; 10.503-8.206", "5.503-8.306; 10.503-8.906"),
        # Imagem 6, Legenda 6
        ("10.703-5.206; 15.703-8.206", "10.703-8.306; 15.703-8.906"),
    ]
    # Coordenadas do título principal
    left_tit, top_tit, width_tit, height_tit = converter_coordenadas_16_9("1-0.3; 15-1.2")
    adicionar_caixa_texto(slide, titulo, left_tit, top_tit, width_tit, height_tit, font_size=36)

    num_legendas = len(legendas)
    num_imagens = len(imagens_bytes)

    for i, (coord_img, coord_txt) in enumerate(coords):
        # Adiciona a imagem se disponível
        if i < num_imagens:
            left_img, top_img, width_img, height_img = converter_coordenadas_16_9(coord_img)
            adicionar_imagem(slide, imagens_bytes[i], left_img, top_img, width_img, height_img)
        else:
            # Opcional: Adicionar um placeholder ou aviso se a imagem estiver em falta
            print(f"Aviso: Imagem {i+1} em falta para o layout img6.")

        # Adiciona a legenda se disponível
        left_txt, top_txt, width_txt, height_txt = converter_coordenadas_16_9(coord_txt)
        legenda_texto = legendas[i] if i < num_legendas else "" # Usa legenda ou string vazia
        adicionar_caixa_texto(slide, legenda_texto, left_txt, top_txt, width_txt, height_txt, font_size=12) # Tamanho de fonte menor para legendas

    if num_legendas < 6:
        print(f"Aviso: Fornecidas {num_legendas} legendas para o layout img6, esperadas 6.")
    if num_imagens < 6:
         print(f"Aviso: Fornecidas {num_imagens} imagens para o layout img6, esperadas 6.")


layout_functions = {
    "img left custom": layout_img_left_custom,
    "img top custom": layout_img_top_custom,
    "img right custom": layout_img_right_custom,
    "img2": layout_img2,
    "img6": layout_img6
}

# --- Função Principal Adaptada ---

def gerar_apresentacao_gui():
    theme_path = theme_var.get()
    image_paths = list(image_listbox.get(0, tk.END))
    comandos_str = text_input.get("1.0", tk.END).strip()

    if not theme_path:
        messagebox.showwarning("Ficheiro em Falta", "Por favor, selecione um ficheiro de tema (.pptx).")
        return
    # Não é estritamente necessário ter imagens se os comandos não as usarem
    # if not image_paths:
    #     messagebox.showwarning("Ficheiros em Falta", "Por favor, selecione pelo menos uma imagem.")
    #     return
    if not comandos_str:
        messagebox.showwarning("Comandos em Falta", "Por favor, insira os comandos para os slides.")
        return

    try:
        # Ler o tema para obter bytes e contar slides iniciais
        initial_slides_count = 0
        theme_bytes = None
        try:
            with open(theme_path, 'rb') as f_theme:
                theme_bytes = f_theme.read()
            # Abrir uma instância separada para contar slides sem afetar a principal
            prs_check = Presentation(BytesIO(theme_bytes))
            initial_slides_count = len(prs_check.slides)
            print(f"Tema original contém {initial_slides_count} slide(s).")
        except Exception as e_theme:
             messagebox.showerror("Erro ao Ler Tema", f"Não foi possível ler ou analisar o ficheiro de tema: {e_theme}")
             return

        # Criar a apresentação principal a partir dos bytes do tema
        prs = Presentation(BytesIO(theme_bytes))

        # Carregar imagens em memória
        imagens_bytes = []
        for img_path in image_paths:
            try:
                with open(img_path, 'rb') as f_img:
                    imagens_bytes.append(f_img.read())
            except Exception as e_img:
                 messagebox.showwarning("Erro ao Ler Imagem", f"Não foi possível ler a imagem: {img_path}\nErro: {e_img}\nEsta imagem será ignorada.")

        comandos = comandos_str.split("\n")
        imagem_idx = 0
        slides_adicionados = 0
        comandos_validos = [] # Para contar comandos que realmente geram slides

        # Processar comandos e adicionar slides
        for comando_idx, comando in enumerate(comandos):
            comando = comando.strip()
            if not comando: continue # Ignorar linhas vazias

            partes = comando.split("|")
            if len(partes) < 2: # Layout e Título são o mínimo
                print(f"Aviso: Comando inválido (partes insuficientes): {comando}")
                status_var.set(f"Erro no comando {comando_idx+1}: partes insuficientes")
                continue

            layout = partes[0].strip().lower()
            titulo = partes[1].strip()
            # O texto/legendas podem estar em partes[2] ou mais além para img6
            texto_parte = partes[2].strip() if len(partes) > 2 else ""

            if layout not in layout_functions:
                print(f"Aviso: Layout '{layout}' não reconhecido, comando ignorado: {comando}")
                status_var.set(f"Erro no comando {comando_idx+1}: layout desconhecido")
                continue

            # Tentar obter o layout em branco (índice 6 é comum, mas não garantido)
            # Usar o primeiro layout disponível como fallback se o 6 não existir
            slide_layout_index = 6
            if slide_layout_index >= len(prs.slide_layouts):
                slide_layout_index = 0 # Tenta usar o primeiro layout
                if not prs.slide_layouts:
                     messagebox.showerror("Erro de Tema", "O tema PPTX não contém nenhuns layouts de slide.")
                     return
                print(f"Aviso: Layout em branco (índice 6) não encontrado. Usando layout {slide_layout_index} como base.")

            try:
                 slide_layout = prs.slide_layouts[slide_layout_index]
            except IndexError:
                 # Isto não deveria acontecer se a verificação anterior funcionou
                 messagebox.showerror("Erro de Tema", f"Não foi possível aceder ao layout de slide índice {slide_layout_index}.")
                 return

            # Adicionar o novo slide
            slide = prs.slides.add_slide(slide_layout)
            slides_adicionados += 1
            comandos_validos.append(comando) # Adiciona à lista de comandos que geraram slide
            layout_func = layout_functions[layout]

            try:
                if layout == "img6":
                    imagens_necessarias = 6
                    # As legendas são todas as partes a partir do índice 2
                    legendas = [p.strip() for p in partes[2:]]
                    num_legendas_fornecidas = len(legendas)

                    if num_legendas_fornecidas < imagens_necessarias:
                         print(f"Aviso: Comando '{comando}' para 'img6' forneceu apenas {num_legendas_fornecidas} legendas, esperadas {imagens_necessarias}. Legendas em falta ficarão em branco.")
                         # Preenche com strings vazias para evitar erros
                         legendas.extend(["" for _ in range(imagens_necessarias - num_legendas_fornecidas)])

                    if imagem_idx + imagens_necessarias > len(imagens_bytes):
                        msg_erro_img = f"Layout 'img6' requer {imagens_necessarias} imagens, mas apenas {len(imagens_bytes) - imagem_idx} estão disponíveis a partir do índice atual. Slide pode estar incompleto."
                        messagebox.showwarning("Imagens Insuficientes", msg_erro_img + f"\nComando: {comando}")
                        print(msg_erro_img)
                        # Continua a execução, mas o layout ficará com imagens em falta
                        imagens_usadas = imagens_bytes[imagem_idx : ] # Usa as que restam
                        imagem_idx = len(imagens_bytes) # Esgota as imagens
                    else:
                        imagens_usadas = imagens_bytes[imagem_idx : imagem_idx + imagens_necessarias]
                        imagem_idx += imagens_necessarias

                    # Chama a função de layout com as legendas como argumentos separados
                    layout_func(slide, titulo, *legendas[:imagens_necessarias], imagens_bytes=imagens_usadas, slide_width=prs.slide_width, slide_height=prs.slide_height)

                else: # Outros layouts
                    # Usa '///' como separador para múltiplos blocos de texto
                    textos = [t.strip() for t in texto_parte.split("///")]
                    texto_combinado = "\n".join(textos)

                    # Verifica se precisa de imagem
                    imagem_necessaria = layout in ["img left custom", "img top custom", "img right custom", "img2"]
                    imagem = None
                    if imagem_necessaria:
                        if imagem_idx >= len(imagens_bytes):
                            msg_erro_img = f"Não há mais imagens disponíveis para o comando: {comando}"
                            messagebox.showwarning("Imagens Insuficientes", msg_erro_img)
                            print(msg_erro_img)
                            # Continua, mas sem imagem
                            imagem_bytes_usada = BytesIO() # Imagem vazia para evitar erro? Ou passar None? Melhor não chamar a função.
                            # Não chamar layout_func se a imagem é essencial? Ou deixar a função lidar com imagem None?
                            # Por agora, vamos passar bytes vazios e deixar a função adicionar_imagem falhar graciosamente.
                            layout_func(slide, titulo, texto_combinado, BytesIO(), prs.slide_width, prs.slide_height)
                        else:
                            imagem_bytes_usada = imagens_bytes[imagem_idx]
                            layout_func(slide, titulo, texto_combinado, imagem_bytes_usada, prs.slide_width, prs.slide_height)
                            imagem_idx += 1 # Avança para a próxima imagem
                    else:
                         # Layouts que não usam imagem (se existissem)
                         # layout_func(slide, titulo, texto_combinado, None, prs.slide_width, prs.slide_height)
                         pass # Atualmente todos os layouts usam imagem

            except Exception as e_layout:
                 err_msg = f"Erro ao aplicar layout '{layout}' para o comando '{comando}': {e_layout}"
                 messagebox.showerror("Erro no Layout", err_msg)
                 print(err_msg)
                 # O slide foi adicionado, mas falhou. Pode ser melhor removê-lo?
                 # A remoção de slides específicos é complexa. Deixar o slide potencialmente quebrado.
                 status_var.set(f"Erro no comando {comando_idx+1}: {e_layout}")
                 continue # Pula para o próximo comando

        # --- Remoção de Slides Excedentes --- 
        slides_realmente_criados = len(comandos_validos)
        total_slides_final = len(prs.slides)
        slides_a_remover = total_slides_final - slides_realmente_criados

        print(f"Total de slides no final: {total_slides_final}")
        print(f"Slides que deveriam ter sido criados: {slides_realmente_criados}")
        print(f"Slides a remover (presumivelmente do tema original no início): {slides_a_remover}")

        if slides_a_remover > 0 and slides_a_remover == initial_slides_count:
            print(f"Tentando remover {slides_a_remover} slides do início...")
            try:
                xml_slides = prs.slides._sldIdLst
                rIds_to_remove = [xml_slides[i].rId for i in range(slides_a_remover)]

                # Remove os elementos sldId do XML
                for _ in range(slides_a_remover):
                    xml_slides.remove(xml_slides[0])

                # Remove as relações correspondentes da parte da apresentação
                pres_part = prs.part
                rels_dropped = 0
                for rId in rIds_to_remove:
                    try:
                        pres_part.drop_rel(rId)
                        rels_dropped += 1
                    except KeyError:
                        print(f"Aviso: Relação {rId} não encontrada para remoção.")
                print(f"Removidos {len(rIds_to_remove)} elementos de slide e {rels_dropped} relações.")
                status_var.set(f"Gerado com {slides_realmente_criados} slides. {slides_a_remover} slides do tema removidos.")

            except Exception as e_rem:
                print(f"Erro ao tentar remover slides excedentes: {e_rem}")
                messagebox.showwarning("Erro na Remoção", "Ocorreu um erro ao tentar remover slides excedentes do tema. Eles podem permanecer no ficheiro final.")
                status_var.set(f"Gerado com {slides_realmente_criados} slides. Falha ao remover slides do tema.")
        elif slides_a_remover > 0:
             print(f"Aviso: Número de slides a remover ({slides_a_remover}) não corresponde ao número inicial de slides do tema ({initial_slides_count}). Remoção automática ignorada.")
             messagebox.showwarning("Remoção Ignorada", "Não foi possível determinar com segurança quais slides remover. Os slides do tema original podem permanecer.")
             status_var.set(f"Gerado com {slides_realmente_criados} slides. Slides do tema não removidos.")
        else:
            status_var.set(f"Gerado com {slides_realmente_criados} slides.")
        # --- Fim da Remoção --- 

        # Guardar a apresentação
        output_path = filedialog.asksaveasfilename(
            defaultextension=".pptx",
            filetypes=[("PowerPoint Presentations", "*.pptx"), ("All Files", "*.*")],
            title="Guardar Apresentação Como..."
        )

        if output_path:
            try:
                prs.save(output_path)
                messagebox.showinfo("Sucesso", f"Apresentação guardada com sucesso em:\n{output_path}")
                # Atualiza status final apenas se salvou com sucesso
                final_status = status_var.get() # Pega o status da remoção/geração
                status_var.set(f"{final_status} Guardado em {os.path.basename(output_path)}")
            except Exception as e_save:
                 messagebox.showerror("Erro ao Guardar", f"Não foi possível guardar a apresentação:\n{e_save}")
                 status_var.set("Erro ao guardar o ficheiro.")
        else:
            status_var.set("Geração cancelada pelo utilizador.")

    except Exception as e:
        messagebox.showerror("Erro Geral", f"Ocorreu um erro inesperado durante a geração: {e}")
        status_var.set("Erro durante a geração.")
        import traceback
        print(traceback.format_exc()) # Imprime traceback para debug

# --- Funções da GUI ---

def selecionar_tema():
    filepath = filedialog.askopenfilename(
        title="Selecionar Ficheiro de Tema",
        filetypes=[("PowerPoint Presentations", "*.pptx")]
    )
    if filepath:
        theme_var.set(filepath)
        status_var.set(f"Tema: {os.path.basename(filepath)}")

def adicionar_imagens():
    filepaths = filedialog.askopenfilenames(
        title="Selecionar Imagens",
        filetypes=[("Image Files", "*.png *.jpg *.jpeg")]
    )
    if filepaths:
        current_list = list(image_listbox.get(0, tk.END))
        added_count = 0
        for fp in filepaths:
            if fp not in current_list:
                image_listbox.insert(tk.END, fp)
                added_count += 1
        if added_count > 0:
             update_image_status()
        status_var.set(f"{added_count} imagem(ns) adicionada(s). Total: {image_listbox.size()}")

def remover_imagem_selecionada():
    selected_indices = image_listbox.curselection()
    if not selected_indices:
        messagebox.showinfo("Nenhuma Seleção", "Selecione uma ou mais imagens para remover.")
        return
    # Remover da lista de trás para a frente para evitar problemas de índice
    removed_count = 0
    for i in reversed(selected_indices):
        image_listbox.delete(i)
        removed_count += 1
    if removed_count > 0:
        update_image_status()
    status_var.set(f"{removed_count} imagem(ns) removida(s). Total: {image_listbox.size()}")

def mover_imagem_cima():
    selected_indices = image_listbox.curselection()
    if not selected_indices:
        messagebox.showinfo("Nenhuma Seleção", "Selecione uma imagem para mover.")
        return
    # Mover apenas o primeiro selecionado para cima
    idx = selected_indices[0]
    if idx > 0:
        text = image_listbox.get(idx)
        image_listbox.delete(idx)
        image_listbox.insert(idx - 1, text)
        image_listbox.selection_set(idx - 1)
        image_listbox.activate(idx - 1)

def mover_imagem_baixo():
    selected_indices = image_listbox.curselection()
    if not selected_indices:
        messagebox.showinfo("Nenhuma Seleção", "Selecione uma imagem para mover.")
        return
    # Mover apenas o primeiro selecionado para baixo
    idx = selected_indices[0]
    if idx < image_listbox.size() - 1:
        text = image_listbox.get(idx)
        image_listbox.delete(idx)
        image_listbox.insert(idx + 1, text)
        image_listbox.selection_set(idx + 1)
        image_listbox.activate(idx + 1)

def update_image_status():
    count = image_listbox.size()
    image_count_var.set(f"Imagens Carregadas: {count}")
    # Também atualiza a status bar geral, mas foca no contador dedicado
    # status_var.set(f"{count} imagem(ns) selecionada(s).")

def reset_all_fields():
    if messagebox.askyesno("Confirmar Limpeza", "Tem a certeza que deseja limpar todos os campos (tema, imagens, comandos)?"):
        theme_var.set("")
        image_listbox.delete(0, tk.END)
        text_input.delete("1.0", tk.END)
        update_image_status() # Atualiza contador para 0
        status_var.set("Campos limpos. Pronto.")

# --- Configuração da GUI --- 
root = tk.Tk()
root.title("Gerador de Apresentações PPTX v2")

# Variáveis Tkinter
theme_var = tk.StringVar()
status_var = tk.StringVar(value="Pronto.")
image_count_var = tk.StringVar(value="Imagens Carregadas: 0") # Variável para contador de imagens

# Frame Principal
main_frame = ttk.Frame(root, padding="10")
main_frame.grid(row=0, column=0, sticky=(tk.W, tk.E, tk.N, tk.S))
root.columnconfigure(0, weight=1)
root.rowconfigure(0, weight=1)

# --- Linha 0: Botão Limpar ---
# Colocado numa linha superior para fácil acesso
clear_button = ttk.Button(main_frame, text="Limpar Tudo", command=reset_all_fields)
clear_button.grid(row=0, column=2, sticky=tk.E, pady=(0, 5), padx=5)

# --- Linha 1: Seleção de Tema ---
theme_label = ttk.Label(main_frame, text="Ficheiro Tema (.pptx):")
theme_label.grid(row=1, column=0, sticky=tk.W, pady=2)
theme_entry = ttk.Entry(main_frame, textvariable=theme_var, width=50, state='readonly')
theme_entry.grid(row=1, column=1, sticky=(tk.W, tk.E), padx=5, pady=2)
theme_button = ttk.Button(main_frame, text="Selecionar...", command=selecionar_tema)
theme_button.grid(row=1, column=2, sticky=tk.E, pady=2, padx=5)

# --- Linha 2: Seleção de Imagens ---
image_controls_frame = ttk.Frame(main_frame)
image_controls_frame.grid(row=2, column=0, columnspan=3, sticky=(tk.W, tk.E))

image_label = ttk.Label(image_controls_frame, text="Ficheiros de Imagem (.png, .jpg):")
image_label.pack(side=tk.LEFT, anchor=tk.W)

# Contador de imagens dedicado
image_count_label = ttk.Label(image_controls_frame, textvariable=image_count_var)
image_count_label.pack(side=tk.RIGHT, anchor=tk.E, padx=5)

# --- Linha 3: Lista de Imagens e Botões ---
image_frame = ttk.Frame(main_frame)
image_frame.grid(row=3, column=0, columnspan=3, sticky=(tk.W, tk.E, tk.N, tk.S), pady=5)
image_frame.columnconfigure(0, weight=1)
image_frame.rowconfigure(0, weight=1)

image_listbox = tk.Listbox(image_frame, selectmode=tk.EXTENDED, height=8)
image_listbox.grid(row=0, column=0, sticky=(tk.W, tk.E, tk.N, tk.S), padx=(0, 5))

img_scrollbar_y = ttk.Scrollbar(image_frame, orient=tk.VERTICAL, command=image_listbox.yview)
img_scrollbar_y.grid(row=0, column=1, sticky=(tk.N, tk.S))
image_listbox['yscrollcommand'] = img_scrollbar_y.set

img_scrollbar_x = ttk.Scrollbar(image_frame, orient=tk.HORIZONTAL, command=image_listbox.xview)
img_scrollbar_x.grid(row=1, column=0, sticky=(tk.E, tk.W))
image_listbox['xscrollcommand'] = img_scrollbar_x.set

img_button_frame = ttk.Frame(image_frame)
img_button_frame.grid(row=0, column=2, sticky=(tk.N, tk.S), padx=(5, 0))

add_img_button = ttk.Button(img_button_frame, text="Adicionar...", command=adicionar_imagens)
add_img_button.pack(pady=2, fill=tk.X)
remove_img_button = ttk.Button(img_button_frame, text="Remover", command=remover_imagem_selecionada)
remove_img_button.pack(pady=2, fill=tk.X)
move_up_button = ttk.Button(img_button_frame, text="Cima", command=mover_imagem_cima)
move_up_button.pack(pady=2, fill=tk.X)
move_down_button = ttk.Button(img_button_frame, text="Baixo", command=mover_imagem_baixo)
move_down_button.pack(pady=2, fill=tk.X)

# --- Linha 4: Comandos ---
cmd_label = ttk.Label(main_frame, text="Comandos (um por linha: layout|título|texto_ou_legenda1|legenda2...):")
cmd_label.grid(row=4, column=0, columnspan=3, sticky=tk.W, pady=(10, 2))

text_input = tk.Text(main_frame, height=10, wrap=tk.WORD)
text_input.grid(row=5, column=0, columnspan=3, sticky=(tk.W, tk.E, tk.N, tk.S), pady=5)

text_scrollbar = ttk.Scrollbar(main_frame, orient=tk.VERTICAL, command=text_input.yview)
text_scrollbar.grid(row=5, column=3, sticky=(tk.N, tk.S), pady=5)
text_input['yscrollcommand'] = text_scrollbar.set

# --- Linha 5: Botão Gerar ---
run_button = ttk.Button(main_frame, text="Gerar Apresentação", command=gerar_apresentacao_gui)
run_button.grid(row=6, column=0, columnspan=3, pady=10)

# --- Barra de Status --- 
status_bar = ttk.Label(root, textvariable=status_var, relief=tk.SUNKEN, anchor=tk.W, padding=5)
status_bar.grid(row=1, column=0, sticky=(tk.W, tk.E))

# Configurar pesos das colunas/linhas para redimensionamento
main_frame.columnconfigure(1, weight=1)
main_frame.rowconfigure(3, weight=1) # Frame da lista de imagens
main_frame.rowconfigure(5, weight=2) # Área de texto dos comandos

# Inicializar o contador de imagens
update_image_status()

# Iniciar loop da GUI
root.mainloop()

